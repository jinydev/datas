import os
import glob
from pptx import Presentation
from pypdf import PdfReader

def convert_pptx_to_md(filepath):
    try:
        prs = Presentation(filepath)
        md_content = f"# {os.path.basename(filepath)}\n\n"
        
        for i, slide in enumerate(prs.slides):
            md_content += f"## Slide {i + 1}\n\n"
            
            # Extract title
            if slide.shapes.title:
                md_content += f"### {slide.shapes.title.text}\n\n"
            
            # Extract text from other shapes
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape == slide.shapes.title:
                    continue
                
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        md_content += f"- {text}\n"
            
            md_content += "\n---\n\n"
            
        return md_content
    except Exception as e:
        print(f"Error converting PPTX {filepath}: {e}")
        return None

def convert_pdf_to_md(filepath):
    try:
        reader = PdfReader(filepath)
        md_content = f"# {os.path.basename(filepath)}\n\n"
        
        for i, page in enumerate(reader.pages):
            md_content += f"## Page {i + 1}\n\n"
            text = page.extract_text()
            if text:
                md_content += f"{text}\n\n"
            md_content += "\n---\n\n"
            
        return md_content
    except Exception as e:
        print(f"Error converting PDF {filepath}: {e}")
        return None

def main():
    # Target directory relative to this script
    target_dir = "."
    files = sorted(os.listdir(target_dir))
    
    for filename in files:
        filepath = os.path.join(target_dir, filename)
        if not os.path.isfile(filepath):
            continue
            
        name, ext = os.path.splitext(filename)
        ext = ext.lower()
        
        md_content = None
        if ext == '.pptx':
            print(f"Converting {filename}...")
            md_content = convert_pptx_to_md(filepath)
        elif ext == '.pdf':
            print(f"Converting {filename}...")
            md_content = convert_pdf_to_md(filepath)
            
        if md_content:
            md_filename = f"{name}.md"
            with open(md_filename, 'w', encoding='utf-8') as f:
                f.write(md_content)
            print(f"Saved {md_filename}")

if __name__ == "__main__":
    main()
